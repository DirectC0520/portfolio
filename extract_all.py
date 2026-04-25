#!/usr/bin/env python3
"""
Extract all content from uploaded files and organize for web portfolio.
"""

import os
import json
import shutil
import glob
from pathlib import Path
from collections import Counter

# ============================================================
# 1. Create directory structure
# ============================================================
dirs = [
    './research/pptx_content/',
    './research/resume_content/',
    './research/script_content/',
    './final/images/',
    './final/css/',
    './final/js/',
    './assets/screenshots/',
    './assets/style_references/',
    './assets/originals/',
    './temp/',
]
for d in dirs:
    os.makedirs(d, exist_ok=True)
print("✅ Directory structure created")

# ============================================================
# 2. Extract content from PPTX
# ============================================================
from pptx import Presentation
from pptx.util import Inches, Pt

pptx_path = './past works ShiHao.pptx'
prs = Presentation(pptx_path)

slides_data = []
image_count = 0

for slide_idx, slide in enumerate(prs.slides, 1):
    slide_info = {
        'slide_number': slide_idx,
        'texts': [],
        'images': []
    }
    
    # Extract text from all shapes
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    slide_info['texts'].append(text)
        
        # Check for tables
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    slide_info['texts'].append(' | '.join(row_texts))
    
    # Extract images
    img_idx = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture type
            img_idx += 1
            image = shape.image
            ext = image.content_type.split('/')[-1]
            if ext == 'jpeg':
                ext = 'jpg'
            elif ext == 'x-emf' or ext == 'x-wmf':
                ext = ext.replace('x-', '')
            
            img_name = f'slide{slide_idx}_image{img_idx}.{ext}'
            img_bytes = image.blob
            
            # Save to research
            research_path = f'./research/pptx_content/{img_name}'
            with open(research_path, 'wb') as f:
                f.write(img_bytes)
            
            # Copy to final/images (only web-ready formats)
            if ext in ('jpg', 'png', 'gif', 'webp', 'jpeg'):
                final_path = f'./final/images/{img_name}'
                with open(final_path, 'wb') as f:
                    f.write(img_bytes)
                slide_info['images'].append({
                    'filename': img_name,
                    'research_path': research_path,
                    'final_path': final_path,
                    'content_type': image.content_type
                })
            else:
                slide_info['images'].append({
                    'filename': img_name,
                    'research_path': research_path,
                    'final_path': None,
                    'content_type': image.content_type,
                    'note': f'Non-web format ({ext}), saved to research only'
                })
            
            image_count += 1
    
    slides_data.append(slide_info)

# Save slides text JSON
with open('./research/pptx_content/slides_text.json', 'w', encoding='utf-8') as f:
    json.dump(slides_data, f, indent=2, ensure_ascii=False)

print(f"✅ PPTX extracted: {len(slides_data)} slides, {image_count} images")
for s in slides_data:
    print(f"   Slide {s['slide_number']}: {len(s['texts'])} text blocks, {len(s['images'])} images")

# ============================================================
# 3. Extract text from Resume PDF
# ============================================================
import pdfplumber

resume_path = './Resume Shihao Lai.pdf'
resume_text = ''
try:
    with pdfplumber.open(resume_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                resume_text += page_text + '\n\n'
    
    with open('./research/resume_content/resume_text.txt', 'w', encoding='utf-8') as f:
        f.write(resume_text)
    print(f"✅ Resume extracted: {len(resume_text)} characters, {len(resume_text.splitlines())} lines")
except Exception as e:
    print(f"⚠️ Resume extraction error: {e}")
    resume_text = f"Error extracting: {e}"

# ============================================================
# 4. Extract text from script.pdf
# ============================================================
script_path = './script.pdf'
script_text = ''
try:
    with pdfplumber.open(script_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                script_text += page_text + '\n\n'
    
    with open('./research/script_content/script_text.txt', 'w', encoding='utf-8') as f:
        f.write(script_text)
    print(f"✅ Script extracted: {len(script_text)} characters, {len(script_text.splitlines())} lines")
except Exception as e:
    print(f"⚠️ Script extraction error: {e}")
    script_text = f"Error extracting: {e}"

# ============================================================
# 5. Organize screenshot images
# ============================================================
screenshots = sorted(glob.glob('./Screenshot *.jpg'))
screenshot_mapping = []

for idx, src in enumerate(screenshots, 1):
    original_name = os.path.basename(src)
    clean_name = f'screenshot_{idx:02d}.jpg'
    
    # Copy to assets/screenshots
    shutil.copy2(src, f'./assets/screenshots/{original_name}')
    
    # Copy to final/images with clean name
    shutil.copy2(src, f'./final/images/{clean_name}')
    
    screenshot_mapping.append({
        'original_name': original_name,
        'clean_name': clean_name,
        'assets_path': f'./assets/screenshots/{original_name}',
        'final_path': f'./final/images/{clean_name}'
    })

print(f"✅ Screenshots organized: {len(screenshot_mapping)} files")

# ============================================================
# 6. Organize storyboard
# ============================================================
storyboard_src = './Storyboard Shihao Lai.png'
shutil.copy2(storyboard_src, './assets/originals/Storyboard Shihao Lai.png')
shutil.copy2(storyboard_src, './final/images/storyboard.png')
print("✅ Storyboard organized")

# ============================================================
# 7. Organize style references
# ============================================================
style_refs = []
for i in range(1, 4):
    src = f'./Style{i}.jpg'
    if os.path.exists(src):
        shutil.copy2(src, f'./assets/style_references/Style{i}.jpg')
        style_refs.append(f'Style{i}.jpg')
print(f"✅ Style references organized: {len(style_refs)} files")

# ============================================================
# 8. Organize originals
# ============================================================
originals = [
    './past works ShiHao.pptx',
    './Resume Shihao Lai.pdf',
    './script.pdf',
]
for src in originals:
    if os.path.exists(src):
        shutil.copy2(src, f'./assets/originals/{os.path.basename(src)}')
print("✅ Originals organized")

# ============================================================
# 9. Analyze style reference images
# ============================================================
from PIL import Image
import struct

def get_dominant_colors(img_path, num_colors=5):
    """Get dominant colors from an image."""
    img = Image.open(img_path)
    # Resize for faster processing
    img_small = img.resize((100, 100))
    img_rgb = img_small.convert('RGB')
    pixels = list(img_rgb.getdata())
    
    # Simple color binning
    color_counts = Counter()
    for r, g, b in pixels:
        # Bin to nearest 32
        binned = (r // 32 * 32, g // 32 * 32, b // 32 * 32)
        color_counts[binned] += 1
    
    top_colors = color_counts.most_common(num_colors)
    return [{'rgb': list(c), 'hex': '#{:02x}{:02x}{:02x}'.format(*c), 'percentage': round(count/len(pixels)*100, 1)} for c, count in top_colors]

style_analysis = []
for i in range(1, 4):
    img_path = f'./assets/style_references/Style{i}.jpg'
    if os.path.exists(img_path):
        img = Image.open(img_path)
        file_size = os.path.getsize(img_path)
        dominant = get_dominant_colors(img_path)
        
        analysis = {
            'filename': f'Style{i}.jpg',
            'width': img.width,
            'height': img.height,
            'file_size_bytes': file_size,
            'file_size_kb': round(file_size / 1024, 1),
            'mode': img.mode,
            'format': img.format,
            'dominant_colors': dominant,
        }
        style_analysis.append(analysis)
        print(f"✅ Style{i}.jpg: {img.width}x{img.height}, {analysis['file_size_kb']}KB")
        for c in dominant[:3]:
            print(f"   Color: {c['hex']} ({c['percentage']}%)")

# ============================================================
# 10. Build comprehensive extracted_data.json
# ============================================================

# List all images in final/images/
all_portfolio_images = []
for f in sorted(os.listdir('./final/images/')):
    fpath = f'./final/images/{f}'
    if os.path.isfile(fpath):
        fsize = os.path.getsize(fpath)
        all_portfolio_images.append({
            'filename': f,
            'path': f'images/{f}',  # relative to final/
            'file_size_bytes': fsize,
            'file_size_kb': round(fsize / 1024, 1)
        })

extracted_data = {
    'resume_text': resume_text,
    'script_text': script_text,
    'slides': slides_data,
    'screenshots': screenshot_mapping,
    'storyboard': {
        'original_path': './assets/originals/Storyboard Shihao Lai.png',
        'final_path': './final/images/storyboard.png',
        'web_path': 'images/storyboard.png'
    },
    'style_references': style_analysis,
    'all_portfolio_images': all_portfolio_images
}

with open('./research/extracted_data.json', 'w', encoding='utf-8') as f:
    json.dump(extracted_data, f, indent=2, ensure_ascii=False)

print(f"\n✅ extracted_data.json created with {len(all_portfolio_images)} portfolio images")
print(f"   Resume: {len(resume_text)} chars")
print(f"   Script: {len(script_text)} chars")
print(f"   Slides: {len(slides_data)}")
print(f"   Screenshots: {len(screenshot_mapping)}")
print(f"   Style refs: {len(style_analysis)}")

# Print summary of all final images
print("\n📁 All images in final/images/:")
for img in all_portfolio_images:
    print(f"   {img['filename']} ({img['file_size_kb']}KB)")
